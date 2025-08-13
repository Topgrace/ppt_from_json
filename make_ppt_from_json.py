import sys
import os
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE

# --- 설정값 ---
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# 폰트 및 색상 설정 (요청사항 반영)
KOREAN_FONT = '나눔스퀘어라운드 ExtraBold'
MATH_FONT = 'BT수식M'
TEXT_COLOR = RGBColor(255, 255, 255)      # 흰색으로 변경
BACKGROUND_COLOR = RGBColor(0, 0, 0)         # 검은색으로 변경
SEA_BLUE_COLOR = RGBColor(0, 255, 255)       # (0, 255, 255) 청록색(Cyan)으로 변경

# 제목
TITLE_LEFT = Cm(1.5)
TITLE_TOP = Cm(1)
TITLE_WIDTH = Cm(20)
TITLE_HEIGHT = Cm(1.5)
TITLE_FONT_SIZE = Pt(32)
TITLE_FONT_BOLD = True

# 문제 번호
PROBLEM_NUM_LEFT = Cm(1.5)
PROBLEM_NUM_TOP = Cm(3.5)
PROBLEM_NUM_WIDTH = Cm(1.5)
PROBLEM_NUM_HEIGHT = Cm(1.5)
PROBLEM_NUM_FONT_SIZE = Pt(32)

# 질문
QUESTION_LEFT = Cm(3.5)
QUESTION_TOP = Cm(3.5)
QUESTION_WIDTH = Cm(25)
QUESTION_MIN_HEIGHT = Cm(2.0)  # 최소 높이로 변경
QUESTION_FONT_SIZE = Pt(28)

# 선택지
CHOICES_LEFT = Cm(3.5)
CHOICES_WIDTH = Cm(30)
CHOICES_FONT_SIZE = Pt(24)
CHOICE_BOX_HEIGHT = Inches(0.5)
CHOICE_VERTICAL_SPACING = Inches(0.1)

# 이미지
IMAGE_LEFT = Cm(4)
IMAGE_MAX_HEIGHT = Cm(10)
IMAGE_MAX_WIDTH = Cm(15)


def apply_formatting(paragraph, text, default_font=KOREAN_FONT, math_font=MATH_FONT, font_color=TEXT_COLOR, font_size=None, bold=False):
    paragraph.text = ""
    if not text: return

    underline_words = ["잘못", "않은", "않는"]
    pattern = f'({"|".join(underline_words)})'
    parts = re.split(pattern, text)

    for part in parts:
        if not part: continue
        is_underline_part = part in underline_words
        
        segments = re.split(r'([a-zA-Z0-9\+\-\=°\(\):]+)', part)

        for segment_text in segments:
            if not segment_text: continue
            
            is_math = bool(re.match(r'^[a-zA-Z0-9\+\-\=°\(\):]+$', segment_text))
            
            run = paragraph.add_run()
            run.text = segment_text
            font = run.font
            font.color.rgb = font_color
            
            if is_math:
                font.name = math_font
            else:
                font.name = default_font

            if font_size:
                font.size = font_size
            
            font.bold = bold
            if is_underline_part:
                font.underline = True

def create_ppt_from_problems(data, output_filename):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]

    print(f"총 {len(data)}개의 슬라이드를 생성합니다...")

    for i, problem in enumerate(data):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR

        slide_title = problem.get("slide_title", "")
        title_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        p = title_box.text_frame.paragraphs[0]
        # 제목 색상을 SEA_BLUE_COLOR로 적용
        apply_formatting(p, slide_title, font_color=SEA_BLUE_COLOR, font_size=TITLE_FONT_SIZE, bold=TITLE_FONT_BOLD)

        num_box = slide.shapes.add_textbox(PROBLEM_NUM_LEFT, PROBLEM_NUM_TOP, PROBLEM_NUM_WIDTH, PROBLEM_NUM_HEIGHT)
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        num_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        apply_formatting(p, problem.get("number", ""), font_color=SEA_BLUE_COLOR, font_size=PROBLEM_NUM_FONT_SIZE, bold=True)

        question_text = problem.get("question", "")
        question_box = slide.shapes.add_textbox(QUESTION_LEFT, QUESTION_TOP, QUESTION_WIDTH, QUESTION_MIN_HEIGHT)
        tf = question_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        apply_formatting(p, question_text, font_size=QUESTION_FONT_SIZE)

        # 텍스트에 맞게 자동 크기 조정 설정
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        next_element_top = question_box.top + question_box.height + Cm(0.5)

        if problem.get("choices"):
            choice_top = next_element_top
            for choice in problem["choices"]:
                choice_box = slide.shapes.add_textbox(CHOICES_LEFT, choice_top, CHOICES_WIDTH, CHOICE_BOX_HEIGHT)
                p = choice_box.text_frame.paragraphs[0]
                # 선택지 텍스트 색상은 기본값(흰색) 사용
                apply_formatting(p, choice, font_size=CHOICES_FONT_SIZE)
                choice_top += choice_box.height + CHOICE_VERTICAL_SPACING
            next_element_top = choice_top

        image_path = problem.get("included_picture")
        if image_path and os.path.exists(image_path):
            try:
                pic = slide.shapes.add_picture(image_path, IMAGE_LEFT, next_element_top, height=IMAGE_MAX_HEIGHT)
                
                if pic.width > IMAGE_MAX_WIDTH:
                    ratio = IMAGE_MAX_WIDTH / pic.width
                    pic.width = IMAGE_MAX_WIDTH
                    pic.height = int(pic.height * ratio)

            except Exception as e:
                print(f"경고: 문제 {problem.get('number')}의 이미지 '{image_path}'를 삽입하지 못했습니다. ({e})")
        elif image_path:
             print(f"경고: 문제 {problem.get('number')}의 이미지 경로를 찾을 수 없습니다: {image_path}")

        print(f"  - 슬라이드 {i+1} 생성 완료: 문제 {problem.get('number', '')}번")

    prs.save(output_filename)
    print(f"\n성공적으로 '{output_filename}' 파일을 생성했습니다.")


if __name__ == '__main__':
    input_json_path = "problems_data.json"
    output_ppt_filename = "수학_문제_슬라이드_결과.pptx"

    if not os.path.exists(input_json_path):
        print(f"오류: JSON 파일을 찾을 수 없습니다 - {input_json_path}")
        sys.exit(1)

    with open(input_json_path, "r", encoding="utf-8") as f:
        problems_data = json.load(f)

    if problems_data:
        create_ppt_from_problems(problems_data, output_ppt_filename)
    else:
        print("JSON 파일에 문제 데이터가 없습니다.")
